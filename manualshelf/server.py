"""
ManualShelf — PDF & Image Manual Catalog
N4MI | Dan Marshall | Grovetown GA

Flask server, port 8075
SQLite catalog: manualshelf.db
Thumbnails: static/thumbs/
Cross-computer: reads/writes manifest.json to NAS exchange folder
"""

import os, json, shutil, hashlib, datetime, threading, time, subprocess, platform
from pathlib import Path
from flask import Flask, render_template, request, jsonify, send_from_directory, abort, Response
import sqlite3

# ── Config ───────────────────────────────────────────────────────────────────
PORT        = 8075
DB_PATH     = Path(__file__).parent / "manualshelf.db"
THUMB_DIR   = Path(__file__).parent / "static" / "thumbs"
CONFIG_FILE = Path(__file__).parent / "manualshelf_config.json"
THUMB_DIR.mkdir(parents=True, exist_ok=True)

# Supported file extensions
PDF_EXTS   = {".pdf"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}
DOC_EXTS   = {".docx", ".doc"}
SHEET_EXTS = {".xlsx", ".xls", ".csv"}
PPTX_EXTS  = {".pptx", ".ppt"}
TEXT_EXTS  = {".txt", ".md"}
HTML_EXTS  = {".html", ".htm"}
ALL_EXTS   = PDF_EXTS | IMAGE_EXTS | DOC_EXTS | SHEET_EXTS | PPTX_EXTS | TEXT_EXTS | HTML_EXTS

def ext_to_type(ext: str) -> str:
    ext = ext.lower()
    if ext in PDF_EXTS:    return "pdf"
    if ext in IMAGE_EXTS:  return "image"
    if ext in DOC_EXTS:    return "word"
    if ext in SHEET_EXTS:  return "sheet"
    if ext in PPTX_EXTS:   return "pptx"
    if ext in TEXT_EXTS:   return ext.lstrip(".")
    if ext in HTML_EXTS:   return "html"
    return "other"

# ── App ───────────────────────────────────────────────────────────────────────
app = Flask(__name__)

# ── Config helpers ────────────────────────────────────────────────────────────
def load_config():
    if CONFIG_FILE.exists():
        with open(CONFIG_FILE) as f:
            return json.load(f)
    return {
        "scan_folders":   [],
        "exchange_folder": "",
        "computer_name":  platform.node(),
        "manifest_push":  True,
        "manifest_pull":  True,
    }

def save_config(cfg):
    with open(CONFIG_FILE, "w") as f:
        json.dump(cfg, f, indent=2)

# ── Database ──────────────────────────────────────────────────────────────────
def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    with get_db() as db:
        db.executescript("""
            CREATE TABLE IF NOT EXISTS manuals (
                id          INTEGER PRIMARY KEY AUTOINCREMENT,
                file_hash   TEXT UNIQUE NOT NULL,
                file_path   TEXT NOT NULL,
                filename    TEXT NOT NULL,
                title       TEXT NOT NULL,
                description TEXT DEFAULT '',
                file_type   TEXT NOT NULL,
                file_size   INTEGER DEFAULT 0,
                page_count  INTEGER DEFAULT 0,
                thumb_path  TEXT DEFAULT '',
                added_at    TEXT NOT NULL,
                tags        TEXT DEFAULT '',
                make        TEXT DEFAULT '',
                model       TEXT DEFAULT '',
                category    TEXT DEFAULT '',
                favorite    INTEGER DEFAULT 0,
                open_count  INTEGER DEFAULT 0,
                last_opened TEXT DEFAULT ''
            );
            CREATE TABLE IF NOT EXISTS scan_queue (
                id          INTEGER PRIMARY KEY AUTOINCREMENT,
                file_path   TEXT UNIQUE NOT NULL,
                filename    TEXT NOT NULL,
                file_type   TEXT NOT NULL,
                file_size   INTEGER DEFAULT 0,
                scanned_at  TEXT NOT NULL,
                in_catalog  INTEGER DEFAULT 0,
                is_updated  INTEGER DEFAULT 0
            );
            CREATE TABLE IF NOT EXISTS pdf_text (
                manual_id   INTEGER PRIMARY KEY,
                content     TEXT NOT NULL,
                indexed_at  TEXT NOT NULL,
                FOREIGN KEY (manual_id) REFERENCES manuals(id) ON DELETE CASCADE
            );
        """)

init_db()

# ── Full-text PDF indexing ────────────────────────────────────────────────────
def index_pdf_text(manual_id: int, file_path: Path):
    """Extract all text from a PDF and store in pdf_text table."""
    try:
        import fitz
        doc   = fitz.open(str(file_path))
        parts = []
        for page in doc:
            text = page.get_text("text")
            if text.strip():
                parts.append(text)
        doc.close()
        content = "\n".join(parts)
        now = datetime.datetime.now().isoformat()
        with get_db() as db:
            db.execute(
                "INSERT OR REPLACE INTO pdf_text (manual_id, content, indexed_at) VALUES (?,?,?)",
                (manual_id, content, now)
            )
    except Exception as e:
        print(f"PDF text index error for {file_path}: {e}")

def reindex_all_pdfs():
    """Re-index every PDF in the catalog. Called from the Settings button."""
    with get_db() as db:
        rows = db.execute(
            "SELECT id, file_path FROM manuals WHERE file_type='pdf'"
        ).fetchall()
    count = 0
    for row in rows:
        fp = Path(row["file_path"])
        if fp.exists():
            index_pdf_text(row["id"], fp)
            count += 1
    return count

# ── Thumbnail generation ───────────────────────────────────────────────────────
def make_thumb_pdf(file_path: Path, hash_id: str) -> str:
    try:
        import fitz
        doc  = fitz.open(str(file_path))
        page = doc[0]
        mat  = fitz.Matrix(1.5, 1.5)
        pix  = page.get_pixmap(matrix=mat)
        out  = THUMB_DIR / f"{hash_id}.jpg"
        pix.save(str(out))
        doc.close()
        return f"thumbs/{hash_id}.jpg"
    except Exception as e:
        print(f"PDF thumb error for {file_path}: {e}")
        return ""

def make_thumb_image(file_path: Path, hash_id: str) -> str:
    try:
        from PIL import Image
        img = Image.open(str(file_path))
        img.thumbnail((300, 400))
        out = THUMB_DIR / f"{hash_id}.jpg"
        img.convert("RGB").save(str(out), "JPEG", quality=85)
        return f"thumbs/{hash_id}.jpg"
    except Exception as e:
        print(f"Image thumb error for {file_path}: {e}")
        return ""

def make_thumb_pptx(file_path: Path, hash_id: str) -> str:
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw
        prs      = Presentation(str(file_path))
        emu_per_inch = 914400
        dpi      = 96
        w_px     = int(prs.slide_width  / emu_per_inch * dpi)
        h_px     = int(prs.slide_height / emu_per_inch * dpi)
        img      = Image.new("RGB", (w_px, h_px), "#1e3a5f")
        draw     = ImageDraw.Draw(img)
        margin   = 20
        draw.rectangle([margin, margin, w_px-margin, h_px-margin], fill="#ffffff", outline="#cccccc", width=2)
        draw.rectangle([margin, h_px-50, w_px-margin, h_px-margin], fill="#1e3a5f")
        img.thumbnail((300, 400))
        out = THUMB_DIR / f"{hash_id}.jpg"
        img.convert("RGB").save(str(out), "JPEG", quality=85)
        return f"thumbs/{hash_id}.jpg"
    except Exception as e:
        print(f"PPTX thumb error for {file_path}: {e}")
        return ""

def file_hash(file_path: Path) -> str:
    h = hashlib.md5()
    h.update(str(file_path).encode())
    try:
        stat = file_path.stat()
        h.update(str(stat.st_size).encode())
        h.update(str(stat.st_mtime).encode())
    except:
        pass
    return h.hexdigest()[:16]

def get_page_count(file_path: Path) -> int:
    try:
        import fitz
        doc = fitz.open(str(file_path))
        n   = len(doc)
        doc.close()
        return n
    except:
        return 0

def get_slide_count(file_path: Path) -> int:
    try:
        from pptx import Presentation
        return len(Presentation(str(file_path)).slides)
    except:
        return 0

def make_thumbnail(file_path: Path, hash_id: str, ft: str) -> str:
    if ft == "pdf":   return make_thumb_pdf(file_path, hash_id)
    if ft == "image": return make_thumb_image(file_path, hash_id)
    if ft == "pptx":  return make_thumb_pptx(file_path, hash_id)
    return ""

def get_count(file_path: Path, ft: str) -> int:
    if ft == "pdf":  return get_page_count(file_path)
    if ft == "pptx": return get_slide_count(file_path)
    return 0

# ── Folder scanning ────────────────────────────────────────────────────────────
def scan_folder(folder_path: str) -> list:
    results = []
    folder  = Path(folder_path)
    if not folder.exists():
        return results
    for p in folder.rglob("*"):
        if p.suffix.lower() in ALL_EXTS and p.is_file():
            ft = ext_to_type(p.suffix)
            results.append({
                "file_path": str(p),
                "filename":  p.name,
                "file_type": ft,
                "file_size": p.stat().st_size,
            })
    return results

# ── Text preview ──────────────────────────────────────────────────────────────
@app.route("/api/preview/<int:mid>")
def api_preview(mid):
    with get_db() as db:
        row = db.execute("SELECT file_path, file_type FROM manuals WHERE id=?", (mid,)).fetchone()
    if not row:
        return "Not found", 404
    fp = Path(row["file_path"])
    if not fp.exists():
        return "File not found on disk", 404
    try:
        text = fp.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        return f"Could not read file: {e}", 500
    return Response(text, mimetype="text/plain; charset=utf-8")

# ── NAS Manifest ──────────────────────────────────────────────────────────────
def push_manifest():
    cfg = load_config()
    exc = cfg.get("exchange_folder", "").strip()
    if not exc or not cfg.get("manifest_push", True):
        return
    exc_path = Path(exc)
    if not exc_path.exists():
        try:
            exc_path.mkdir(parents=True)
        except:
            return
    computer = cfg.get("computer_name", platform.node())
    with get_db() as db:
        rows = db.execute(
            "SELECT filename, title, description, file_type, file_size, page_count, "
            "tags, make, model, category, added_at FROM manuals ORDER BY title"
        ).fetchall()
    manifest = {
        "computer": computer,
        "updated":  datetime.datetime.now().isoformat(),
        "port":     PORT,
        "count":    len(rows),
        "manuals":  [dict(r) for r in rows],
    }
    out = exc_path / f"manualshelf_{computer}.json"
    try:
        with open(out, "w") as f:
            json.dump(manifest, f, indent=2)
    except Exception as e:
        print(f"Manifest push error: {e}")

def pull_manifests():
    cfg      = load_config()
    exc      = cfg.get("exchange_folder", "").strip()
    computer = cfg.get("computer_name", platform.node())
    if not exc:
        return []
    exc_path = Path(exc)
    if not exc_path.exists():
        return []
    results = []
    for f in exc_path.glob("manualshelf_*.json"):
        if f.stem == f"manualshelf_{computer}":
            continue
        try:
            with open(f) as fh:
                data = json.load(fh)
            results.append(data)
        except:
            pass
    return results

def manifest_loop():
    while True:
        time.sleep(300)
        try:
            push_manifest()
        except:
            pass

threading.Thread(target=manifest_loop, daemon=True).start()

# ── Routes ─────────────────────────────────────────────────────────────────────

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/static/thumbs/<path:filename>")
def serve_thumb(filename):
    return send_from_directory(THUMB_DIR, filename)

# --- Config ---

@app.route("/api/config", methods=["GET"])
def api_get_config():
    return jsonify(load_config())

@app.route("/api/config", methods=["POST"])
def api_save_config():
    cfg  = load_config()
    data = request.json
    cfg.update(data)
    save_config(cfg)
    return jsonify({"ok": True})

# --- Reindex ---

@app.route("/api/reindex", methods=["POST"])
def api_reindex():
    """Re-index all PDF text for full-text search. Run once from Settings."""
    def run():
        count = reindex_all_pdfs()
        print(f"Reindex complete: {count} PDFs indexed")
    threading.Thread(target=run, daemon=True).start()
    with get_db() as db:
        total = db.execute("SELECT COUNT(*) FROM manuals WHERE file_type='pdf'").fetchone()[0]
    return jsonify({"ok": True, "queued": total})

# --- Scan ---

@app.route("/api/scan", methods=["POST"])
def api_scan():
    folder = (request.json or {}).get("folder", "").strip()
    if not folder or not Path(folder).exists():
        return jsonify({"error": "Folder not found"}), 400
    files = scan_folder(folder)
    with get_db() as db:
        existing = {
            r[0]: r[1]
            for r in db.execute("SELECT file_path, file_hash FROM manuals").fetchall()
        }
        now = datetime.datetime.now().isoformat()
        db.execute("DELETE FROM scan_queue")
        for f in files:
            p        = Path(f["file_path"])
            new_hash = file_hash(p)
            in_cat   = 1 if f["file_path"] in existing else 0
            is_upd   = 1 if (f["file_path"] in existing and existing[f["file_path"]] != new_hash) else 0
            try:
                db.execute(
                    "INSERT OR REPLACE INTO scan_queue "
                    "(file_path,filename,file_type,file_size,scanned_at,in_catalog,is_updated) "
                    "VALUES (?,?,?,?,?,?,?)",
                    (f["file_path"], f["filename"], f["file_type"], f["file_size"], now, in_cat, is_upd)
                )
            except:
                pass
    return jsonify({"found": len(files), "files": files})

@app.route("/api/scan_queue")
def api_scan_queue():
    with get_db() as db:
        rows = db.execute("SELECT * FROM scan_queue ORDER BY filename").fetchall()
    return jsonify([dict(r) for r in rows])

# --- Add to catalog ---

@app.route("/api/add", methods=["POST"])
def api_add():
    data      = request.json or {}
    file_path = data.get("file_path", "").strip()
    if not file_path or not Path(file_path).exists():
        return jsonify({"error": "File not found"}), 400
    p   = Path(file_path)
    ft  = ext_to_type(p.suffix)
    hid = file_hash(p)
    now = datetime.datetime.now().isoformat()

    thumb = make_thumbnail(p, hid, ft)
    pages = get_count(p, ft)

    with get_db() as db:
        existing = db.execute(
            "SELECT id, title, description, tags, make, model, category, "
            "favorite, open_count, last_opened, thumb_path "
            "FROM manuals WHERE file_path=?", (str(p),)
        ).fetchone()

        if existing:
            old_thumb = existing["thumb_path"]
            if old_thumb:
                old_tp = THUMB_DIR / Path(old_thumb).name
                old_tp.unlink(missing_ok=True)
            db.execute(
                "UPDATE manuals SET "
                "file_hash=?, filename=?, file_size=?, page_count=?, thumb_path=?, added_at=? "
                "WHERE id=?",
                (hid, p.name, p.stat().st_size, pages, thumb, now, existing["id"])
            )
            db.execute("UPDATE scan_queue SET in_catalog=1, is_updated=0 WHERE file_path=?", (str(p),))
            if ft == "pdf":
                index_pdf_text(existing["id"], p)
            push_manifest()
            return jsonify({"ok": True, "id": existing["id"], "updated": True})

        title = p.stem.replace("_", " ").replace("-", " ")
        cur = db.execute(
            "INSERT OR IGNORE INTO manuals "
            "(file_hash,file_path,filename,title,description,file_type,file_size,"
            "page_count,thumb_path,added_at,tags,make,model,category) "
            "VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?)",
            (hid, str(p), p.name, title, data.get("description",""),
             ft, p.stat().st_size, pages, thumb, now,
             data.get("tags",""), data.get("make",""), data.get("model",""), data.get("category",""))
        )
        new_id = cur.lastrowid
        db.execute("UPDATE scan_queue SET in_catalog=1, is_updated=0 WHERE file_path=?", (str(p),))

    if ft == "pdf":
        index_pdf_text(new_id, p)
    push_manifest()
    return jsonify({"ok": True, "id": new_id})

@app.route("/api/add_batch", methods=["POST"])
def api_add_batch():
    paths   = (request.json or {}).get("paths", [])
    added   = 0
    updated = 0
    errors  = []
    now     = datetime.datetime.now().isoformat()

    for fp in paths:
        p = Path(fp)
        if not p.exists():
            errors.append(fp)
            continue
        ft  = ext_to_type(p.suffix)
        hid = file_hash(p)
        thumb = make_thumbnail(p, hid, ft)
        pages = get_count(p, ft)

        try:
            with get_db() as db:
                existing = db.execute(
                    "SELECT id, thumb_path FROM manuals WHERE file_path=?", (str(p),)
                ).fetchone()

                if existing:
                    old_thumb = existing["thumb_path"]
                    if old_thumb:
                        old_tp = THUMB_DIR / Path(old_thumb).name
                        old_tp.unlink(missing_ok=True)
                    db.execute(
                        "UPDATE manuals SET "
                        "file_hash=?, filename=?, file_size=?, page_count=?, thumb_path=?, added_at=? "
                        "WHERE id=?",
                        (hid, p.name, p.stat().st_size, pages, thumb, now, existing["id"])
                    )
                    if ft == "pdf":
                        index_pdf_text(existing["id"], p)
                    updated += 1
                else:
                    title = p.stem.replace("_", " ").replace("-", " ")
                    cur = db.execute(
                        "INSERT OR IGNORE INTO manuals "
                        "(file_hash,file_path,filename,title,description,file_type,file_size,"
                        "page_count,thumb_path,added_at,tags,make,model,category) "
                        "VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?)",
                        (hid, str(p), p.name, title, "", ft,
                         p.stat().st_size, pages, thumb, now, "", "", "", "")
                    )
                    if ft == "pdf":
                        index_pdf_text(cur.lastrowid, p)
                    added += 1

                db.execute(
                    "UPDATE scan_queue SET in_catalog=1, is_updated=0 WHERE file_path=?", (str(p),)
                )
        except Exception as e:
            errors.append(f"{fp}: {e}")

    push_manifest()
    return jsonify({"added": added, "updated": updated, "errors": errors})

# --- Catalog ---

@app.route("/api/manuals")
def api_manuals():
    q     = request.args.get("q", "").strip().lower()
    tag   = request.args.get("tag", "").strip().lower()
    cat   = request.args.get("category", "").strip()
    ft    = request.args.get("type", "").strip()
    fav   = request.args.get("favorite", "").strip()
    sort  = request.args.get("sort", "title")
    order = request.args.get("order", "asc")

    sort_col = {"title":"title","added":"added_at","name":"filename","opens":"open_count"}.get(sort,"title")
    ord_dir  = "DESC" if order == "desc" else "ASC"

    with get_db() as db:
        rows = db.execute(f"SELECT * FROM manuals ORDER BY {sort_col} {ord_dir}").fetchall()

        # Full-text search: get IDs of PDFs matching query in extracted text
        fts_ids = set()
        if q:
            fts_rows = db.execute(
                "SELECT manual_id FROM pdf_text WHERE LOWER(content) LIKE ?",
                (f"%{q}%",)
            ).fetchall()
            fts_ids = {r[0] for r in fts_rows}

    results = []
    for r in rows:
        d = dict(r)
        if q:
            # Match on metadata fields OR full-text index
            meta_match = (
                q in d["title"].lower() or
                q in d["description"].lower() or
                q in (d["tags"] or "").lower() or
                q in (d["make"] or "").lower() or
                q in (d["model"] or "").lower() or
                q in d["filename"].lower()
            )
            fts_match = d["id"] in fts_ids
            if not meta_match and not fts_match:
                continue
            # Tag results that came from full-text search (UI can show indicator)
            d["fts_match"] = fts_match and not meta_match
        else:
            d["fts_match"] = False

        if tag and tag not in (d["tags"] or "").lower():
            continue
        if cat and cat != d.get("category",""):
            continue
        if ft and ft != d["file_type"]:
            continue
        if fav == "1" and not d["favorite"]:
            continue
        results.append(d)

    return jsonify(results)

@app.route("/api/manuals/<int:mid>", methods=["GET"])
def api_manual_get(mid):
    with get_db() as db:
        row = db.execute("SELECT * FROM manuals WHERE id=?", (mid,)).fetchone()
    if not row:
        return jsonify({"error": "Not found"}), 404
    return jsonify(dict(row))

@app.route("/api/manuals/<int:mid>", methods=["PATCH"])
def api_manual_update(mid):
    data    = request.json or {}
    allowed = ["title","description","tags","make","model","category","favorite"]
    sets    = ", ".join(f"{k}=?" for k in data if k in allowed)
    vals    = [v for k,v in data.items() if k in allowed]
    if not sets:
        return jsonify({"ok": True})
    with get_db() as db:
        db.execute(f"UPDATE manuals SET {sets} WHERE id=?", vals + [mid])
    push_manifest()
    return jsonify({"ok": True})

@app.route("/api/manuals/<int:mid>", methods=["DELETE"])
def api_manual_delete(mid):
    with get_db() as db:
        row = db.execute("SELECT thumb_path FROM manuals WHERE id=?", (mid,)).fetchone()
        if row and row["thumb_path"]:
            tp = THUMB_DIR / Path(row["thumb_path"]).name
            if tp.exists():
                tp.unlink(missing_ok=True)
        db.execute("DELETE FROM manuals WHERE id=?", (mid,))
    push_manifest()
    return jsonify({"ok": True})

# --- Bulk edit ---

@app.route("/api/manuals/bulk", methods=["POST"])
def api_bulk_edit():
    """Apply tags and/or category to multiple manuals at once."""
    data    = request.json or {}
    ids     = data.get("ids", [])
    tags    = data.get("tags", None)       # None = don't touch; "" = clear; "a,b" = set
    category = data.get("category", None)  # same convention
    append_tags = data.get("append_tags", False)  # True = merge with existing tags

    if not ids:
        return jsonify({"ok": True, "updated": 0})

    updated = 0
    with get_db() as db:
        for mid in ids:
            row = db.execute("SELECT tags, category FROM manuals WHERE id=?", (mid,)).fetchone()
            if not row:
                continue
            sets, vals = [], []
            if tags is not None:
                if append_tags and row["tags"]:
                    # Merge: combine existing + new, deduplicate
                    existing_tags = {t.strip() for t in row["tags"].split(",") if t.strip()}
                    new_tags      = {t.strip() for t in tags.split(",") if t.strip()}
                    merged        = ", ".join(sorted(existing_tags | new_tags))
                    sets.append("tags=?"); vals.append(merged)
                else:
                    sets.append("tags=?"); vals.append(tags)
            if category is not None:
                sets.append("category=?"); vals.append(category)
            if sets:
                db.execute(f"UPDATE manuals SET {', '.join(sets)} WHERE id=?", vals + [mid])
                updated += 1

    push_manifest()
    return jsonify({"ok": True, "updated": updated})

# --- Open file ---

@app.route("/api/open/<int:mid>", methods=["POST"])
def api_open(mid):
    with get_db() as db:
        row = db.execute("SELECT file_path FROM manuals WHERE id=?", (mid,)).fetchone()
    if not row:
        return jsonify({"error": "Not found"}), 404
    fp = Path(row["file_path"])
    if not fp.exists():
        return jsonify({"error": f"File not found on disk: {fp}"}), 404
    try:
        if platform.system() == "Windows":
            os.startfile(str(fp))
        elif platform.system() == "Darwin":
            subprocess.Popen(["open", str(fp)])
        else:
            subprocess.Popen(["xdg-open", str(fp)])
        now = datetime.datetime.now().isoformat()
        with get_db() as db:
            db.execute("UPDATE manuals SET open_count=open_count+1, last_opened=? WHERE id=?", (now, mid))
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# --- Tags & categories ---

@app.route("/api/tags")
def api_tags():
    with get_db() as db:
        rows = db.execute("SELECT tags FROM manuals WHERE tags != ''").fetchall()
    tags = set()
    for r in rows:
        for t in r["tags"].split(","):
            t = t.strip()
            if t:
                tags.add(t)
    return jsonify(sorted(tags))

@app.route("/api/categories")
def api_categories():
    with get_db() as db:
        rows = db.execute(
            "SELECT DISTINCT category FROM manuals WHERE category != '' ORDER BY category"
        ).fetchall()
    return jsonify([r["category"] for r in rows])

# --- Remote manifests ---

@app.route("/api/remote")
def api_remote():
    return jsonify(pull_manifests())

# --- Stats ---

@app.route("/api/stats")
def api_stats():
    with get_db() as db:
        total   = db.execute("SELECT COUNT(*) FROM manuals").fetchone()[0]
        pdfs    = db.execute("SELECT COUNT(*) FROM manuals WHERE file_type='pdf'").fetchone()[0]
        images  = db.execute("SELECT COUNT(*) FROM manuals WHERE file_type='image'").fetchone()[0]
        favs    = db.execute("SELECT COUNT(*) FROM manuals WHERE favorite=1").fetchone()[0]
        indexed = db.execute("SELECT COUNT(*) FROM pdf_text").fetchone()[0]
        recents = db.execute(
            "SELECT * FROM manuals WHERE last_opened!='' ORDER BY last_opened DESC LIMIT 5"
        ).fetchall()
    cfg = load_config()
    return jsonify({
        "total": total, "pdfs": pdfs, "images": images, "favorites": favs,
        "indexed": indexed,
        "computer": cfg.get("computer_name", platform.node()),
        "recents": [dict(r) for r in recents],
    })

if __name__ == "__main__":
    print(f"ManualShelf starting on http://localhost:{PORT}")
    push_manifest()
    app.run(host="0.0.0.0", port=PORT, debug=False)
